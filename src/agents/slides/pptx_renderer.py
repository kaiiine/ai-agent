"""PPTX renderer — alternating dark/light professional theme."""
from __future__ import annotations

from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


# ── Palette ───────────────────────────────────────────────────────────────────
def _rgb(r: int, g: int, b: int):
    return RGBColor(r, g, b)

_DARK_BG    = _rgb(0x0d, 0x13, 0x20) if _PPTX_OK else None
_LIGHT_BG   = _rgb(0x11, 0x18, 0x27) if _PPTX_OK else None
_DARK_CARD  = _rgb(0x14, 0x1d, 0x2e) if _PPTX_OK else None
_LIGHT_CARD = _rgb(0xFF, 0xFF, 0xFF) if _PPTX_OK else None
_DARK_TEXT  = _rgb(0xF0, 0xF4, 0xFF) if _PPTX_OK else None
_LIGHT_TEXT = _rgb(0x1a, 0x1f, 0x2e) if _PPTX_OK else None
_MUTED      = _rgb(0x64, 0x74, 0x8b) if _PPTX_OK else None

_ACCENTS_RGB = [
    (0x4a, 0x90, 0xd9),
    (0x00, 0xc8, 0x96),
    (0xf5, 0xa6, 0x23),
    (0x8b, 0x5c, 0xf6),
    (0xe8, 0x54, 0x5a),
    (0x14, 0xb8, 0xa6),
    (0xf9, 0x73, 0x16),
]

_LIGHT_TYPES = {"agenda", "timeline", "stats", "cases", "framework"}

def _accent_rgb(idx: int):
    return _rgb(*_ACCENTS_RGB[idx % len(_ACCENTS_RGB)])

def _is_light(s: dict) -> bool:
    return s.get("type", "content") in _LIGHT_TYPES

# ── Slide dimensions 16:9 ─────────────────────────────────────────────────────
_W = Inches(13.33)
_H = Inches(7.5)


# ── Primitive helpers ─────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, x, y, w, h, fill_color, line_color=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def _txt(slide, text: str, x, y, w, h, *,
         size=14, bold=False, italic=False, color=None,
         align=PP_ALIGN.LEFT, valign=None, wrap=True, font="Calibri"):
    if not text:
        return
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return tb

def _accent_bar(slide, color, right=False):
    x = _W - Pt(5) if right else Inches(0)
    _rect(slide, x, Inches(0), Pt(5), _H, color)


# ── Slide renderers ───────────────────────────────────────────────────────────

def _pptx_title(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(0)
    _accent_bar(sl, ac)

    eyebrow = s.get("eyebrow", "PRÉSENTATION")
    _txt(sl, eyebrow, Inches(0.7), Inches(1.9), Inches(10), Inches(0.4),
         size=10, bold=True, color=ac)
    _txt(sl, s.get("title", ""), Inches(0.7), Inches(2.35), Inches(10), Inches(1.7),
         size=40, bold=True, color=_DARK_TEXT)
    if s.get("subtitle"):
        _txt(sl, s["subtitle"], Inches(0.7), Inches(4.15), Inches(9), Inches(0.8),
             size=16, color=_MUTED)
    meta = "  ·  ".join(x for x in [s.get("author"), s.get("date")] if x)
    if meta:
        _txt(sl, meta, Inches(0.7), Inches(5.2), Inches(9), Inches(0.4), size=12, color=_MUTED)
    if s.get("source"):
        _txt(sl, s["source"], Inches(0.7), Inches(6.9), Inches(12), Inches(0.35),
             size=8, color=_rgb(0x3d, 0x4a, 0x5e))


def _pptx_agenda(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _LIGHT_BG)
    _txt(sl, s.get("title", "Sommaire"), Inches(0.5), Inches(0.3), Inches(12), Inches(0.75),
         size=30, bold=True, color=_DARK_TEXT)

    items = s.get("items", [])
    cols = 3 if len(items) > 3 else max(len(items), 1)
    card_w = Inches(12.5 / cols)
    card_h = Inches(1.5)

    for i, item in enumerate(items[:6]):
        col = i % cols
        row = i // cols
        x = Inches(0.4) + col * (card_w + Inches(0.1))
        y = Inches(1.25) + row * (card_h + Inches(0.15))
        ac = _accent_rgb(i)

        _rect(sl, x, y, card_w - Inches(0.1), card_h, _LIGHT_CARD,
              _rgb(0xe0, 0xe4, 0xee))

        label = item.get("label", item) if isinstance(item, dict) else item
        sub   = item.get("sub", "")   if isinstance(item, dict) else ""

        _txt(sl, str(i + 1).zfill(2), x + Inches(0.15), y + Inches(0.1), Inches(0.6), Inches(0.5),
             size=22, bold=True, color=ac)
        _txt(sl, str(label), x + Inches(0.15), y + Inches(0.62), card_w - Inches(0.35), Inches(0.45),
             size=12, bold=True, color=_LIGHT_TEXT)
        if sub:
            _txt(sl, str(sub), x + Inches(0.15), y + Inches(1.05), card_w - Inches(0.35), Inches(0.35),
                 size=9.5, color=_MUTED)


def _pptx_content(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(idx)
    _accent_bar(sl, ac)

    _txt(sl, s.get("title", ""), Inches(0.7), Inches(0.4), Inches(11), Inches(0.85),
         size=24, bold=True, color=_DARK_TEXT)
    y = Inches(1.4)
    if s.get("body"):
        _txt(sl, s["body"], Inches(0.7), y, Inches(11.7), Inches(1.1),
             size=14, color=_MUTED)
        y += Inches(1.2)
    for b in s.get("bullets", [])[:8]:
        _txt(sl, f"  •  {b}", Inches(0.7), y, Inches(11.7), Inches(0.55),
             size=13, color=_DARK_TEXT)
        y += Inches(0.58)


def _pptx_split(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(idx)
    _accent_bar(sl, ac)

    _txt(sl, s.get("title", ""), Inches(0.7), Inches(0.4), Inches(11), Inches(0.8),
         size=24, bold=True, color=_DARK_TEXT)
    # Divider
    _rect(sl, Inches(6.55), Inches(1.4), Pt(1), Inches(5.8), _DARK_CARD)

    def _col(col: dict, lx: float, col_ac):
        cy = Inches(1.5)
        if col.get("heading"):
            _txt(sl, col["heading"], lx, cy, Inches(5.6), Inches(0.45),
                 size=11, bold=True, color=col_ac)
            cy += Inches(0.55)
        if col.get("body"):
            _txt(sl, col["body"], lx, cy, Inches(5.6), Inches(1.1),
                 size=13, color=_MUTED)
            cy += Inches(1.2)
        for b in col.get("bullets", [])[:6]:
            _txt(sl, f"  •  {b}", lx, cy, Inches(5.6), Inches(0.5),
                 size=12, color=_DARK_TEXT)
            cy += Inches(0.55)

    _col(s.get("left",  {}), Inches(0.7), _accent_rgb(idx))
    _col(s.get("right", {}), Inches(6.8), _accent_rgb(idx + 1))


def _pptx_stats(prs, s: dict, idx: int):
    sl = _blank(prs)
    light = _is_light(s)
    _set_bg(sl, _LIGHT_BG if light else _DARK_BG)

    _txt(sl, s.get("title", ""), Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.75),
         size=28, bold=True, color=_DARK_TEXT,
         align=PP_ALIGN.CENTER)
    if s.get("source"):
        _txt(sl, s["source"], Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.3),
             size=8.5, color=_MUTED, align=PP_ALIGN.CENTER)

    stats = s.get("stats", [])[:6]
    n = len(stats)
    cols = 2 if n == 4 else min(max(n, 1), 3)
    card_w = Inches(12.5 / cols) - Inches(0.15)
    card_h = Inches(3.8)
    x_start = Inches(0.4)
    y_cards = Inches(1.35)

    for i, st in enumerate(stats):
        col = i % cols
        row = i // cols
        ac = _accent_rgb(i)
        x = x_start + col * (card_w + Inches(0.2))
        y = y_cards + row * (card_h + Inches(0.2))

        card_fill = _LIGHT_CARD if light else _DARK_CARD
        _rect(sl, x, y, card_w, card_h, card_fill)
        # Top accent border
        _rect(sl, x, y, card_w, Pt(3.5), ac)

        if st.get("icon"):
            _txt(sl, st["icon"], x, y + Inches(0.4), card_w, Inches(0.6),
                 size=22, align=PP_ALIGN.CENTER)
        _txt(sl, st.get("value", ""), x, y + Inches(1.0), card_w, Inches(1.1),
             size=36, bold=True, color=ac, align=PP_ALIGN.CENTER)
        _txt(sl, st.get("label", ""), x, y + Inches(2.15), card_w, Inches(0.9),
             size=11, color=_MUTED, align=PP_ALIGN.CENTER)
        if st.get("source"):
            _txt(sl, st["source"], x, y + Inches(3.1), card_w, Inches(0.45),
                 size=8, color=_MUTED, align=PP_ALIGN.CENTER)


def _pptx_table(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(idx)

    _txt(sl, s.get("title", ""), Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.72),
         size=26, bold=True, color=_DARK_TEXT)

    cols = s.get("columns", [])
    rows = s.get("rows", [])
    ncols = len(cols)
    dim_w = Inches(2.0)
    col_w = (Inches(12.5) - dim_w) / max(ncols, 1)
    y0 = Inches(1.05)
    row_h = Inches(0.68)

    # Column headers
    for i, col in enumerate(cols):
        label = col.get("label", col) if isinstance(col, dict) else str(col)
        col_ac = _accent_rgb(i)
        _rect(sl, dim_w + i * col_w, y0, col_w - Inches(0.05), row_h - Pt(4), col_ac)
        _txt(sl, label, dim_w + i * col_w, y0, col_w - Inches(0.05), row_h,
             size=11, bold=True,
             color=_LIGHT_TEXT if i == 1 else _DARK_TEXT,
             align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows):
        y = y0 + (r + 1) * row_h + Inches(0.06)
        row_fill = _rgb(0x25, 0x2b, 0x3d) if r % 2 == 0 else _rgb(0x1e, 0x24, 0x35)
        _rect(sl, Inches(0.4), y, Inches(12.5), row_h - Pt(4), row_fill)

        _txt(sl, row.get("dim", ""), Inches(0.5), y, dim_w - Inches(0.1), row_h,
             size=10, bold=True, color=ac)

        for ci, val in enumerate(row.get("values", [])[:ncols]):
            vx = dim_w + ci * col_w + Inches(0.1)
            val_text = "\n".join(val) if isinstance(val, list) else str(val)
            _txt(sl, val_text, vx, y, col_w - Inches(0.2), row_h,
                 size=10,
                 color=_rgb(0xD0, 0xD8, 0xF0) if ci == 0 else _rgb(0xD0, 0xF0, 0xE8))


def _pptx_cases(prs, s: dict, idx: int):
    sl = _blank(prs)
    light = _is_light(s)
    _set_bg(sl, _LIGHT_BG if light else _DARK_BG)

    _txt(sl, s.get("title", ""), Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.72),
         size=28, bold=True, color=_DARK_TEXT)

    cases = s.get("cases", [])[:4]
    card_w = Inches(6.15)
    card_h = Inches(2.85)

    for i, c in enumerate(cases):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * (card_w + Inches(0.2))
        y = Inches(1.05) + row * (card_h + Inches(0.15))
        ac = _rgb(*_ACCENTS_RGB[i % len(_ACCENTS_RGB)])

        card_fill = _LIGHT_CARD if light else _DARK_CARD
        _rect(sl, x, y, card_w, card_h, card_fill)
        # Header
        _rect(sl, x, y, card_w, Inches(0.58), ac)
        company = _e_plain(c.get("company", ""))
        arch    = _e_plain(c.get("arch", ""))
        icon    = c.get("icon", "")

        _txt(sl, f"{icon}  {company}", x + Inches(0.12), y + Inches(0.08),
             Inches(3.5), Inches(0.45), size=12, bold=True,
             color=_LIGHT_TEXT if i % 2 == 0 else _DARK_TEXT)
        _txt(sl, arch, x + card_w - Inches(1.8), y + Inches(0.1), Inches(1.65), Inches(0.38),
             size=9, bold=True,
             color=_LIGHT_TEXT if i % 2 == 0 else _DARK_TEXT,
             align=PP_ALIGN.RIGHT)

        body_y = y + Inches(0.68)
        _txt(sl, _e_plain(c.get("body", "")), x + Inches(0.15), body_y,
             card_w - Inches(0.3), Inches(1.5),
             size=10.5, color=_LIGHT_TEXT if light else _DARK_TEXT)
        if c.get("lesson"):
            _txt(sl, _e_plain(c["lesson"]), x + Inches(0.15), y + card_h - Inches(0.5),
                 card_w - Inches(0.3), Inches(0.38),
                 size=9, italic=True, color=ac)


def _pptx_quote(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _rgb(0x08, 0x06, 0x0f))
    ac = _accent_rgb(idx)
    _txt(sl, "“", Inches(0.7), Inches(0.35), Inches(2), Inches(1.4),
         size=80, bold=True, color=ac)
    _txt(sl, s.get("quote", ""), Inches(0.9), Inches(1.5), Inches(11.4), Inches(3.2),
         size=20, italic=True, align=PP_ALIGN.CENTER, color=_DARK_TEXT)
    author = f"— {s.get('author', '')}"
    if s.get("role"):
        author += f"\n{s['role']}"
    _txt(sl, author, Inches(0.9), Inches(4.9), Inches(11.4), Inches(0.9),
         size=14, color=ac, align=PP_ALIGN.CENTER)


def _pptx_closing(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(0)
    _accent_bar(sl, ac, right=True)

    eyebrow = s.get("eyebrow", "MERCI")
    _txt(sl, eyebrow, Inches(0.7), Inches(1.85), Inches(11), Inches(0.45),
         size=12, bold=True, color=ac)
    _txt(sl, s.get("title", ""), Inches(0.7), Inches(2.35), Inches(11), Inches(1.7),
         size=40, bold=True, color=_DARK_TEXT)
    if s.get("subtitle"):
        _txt(sl, s["subtitle"], Inches(0.7), Inches(4.2), Inches(10), Inches(0.7),
             size=17, color=_MUTED)
    if s.get("cta"):
        _txt(sl, s["cta"], Inches(0.7), Inches(5.1), Inches(6), Inches(0.55),
             size=14, bold=True, color=ac)


def _e_plain(v: Any) -> str:
    return str(v or "")


def _pptx_section(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(idx)

    # Left panel background (42% width)
    left_w = Inches(5.6)
    _rect(sl, Inches(0), Inches(0), left_w, _H, _rgb(0x0a, 0x0f, 0x1a))

    num = str(s.get("num", str(idx).zfill(2)))
    eyebrow = s.get("eyebrow", "PARTIE")

    # Giant ghost number
    _txt(sl, num, Inches(0.1), Inches(-0.3), left_w - Inches(0.1), Inches(5.5),
         size=180, bold=True, color=_rgb(0x28, 0x30, 0x45), align=PP_ALIGN.CENTER, font="Calibri")

    # Eyebrow
    _txt(sl, eyebrow, Inches(6.0), Inches(2.4), Inches(6.8), Inches(0.45),
         size=11, bold=True, color=ac)

    # Vertical bar
    _rect(sl, Inches(6.0), Inches(2.95), Pt(4), Inches(0.9), ac)

    # Title
    _txt(sl, s.get("title", ""), Inches(6.0), Inches(3.0), Inches(6.8), Inches(2.2),
         size=32, bold=True, color=_DARK_TEXT)

    if s.get("subtitle"):
        _txt(sl, s["subtitle"], Inches(6.0), Inches(5.4), Inches(6.8), Inches(0.8),
             size=13, color=_MUTED)


def _pptx_timeline(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _LIGHT_BG)

    _txt(sl, s.get("title", "Programme"), Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.75),
         size=28, bold=True, color=_DARK_TEXT)

    steps = s.get("steps", [])[:6]
    n = max(len(steps), 1)
    step_w = Inches(12.5 / n)
    dot_y  = Inches(2.0)
    dot_r  = Inches(0.45)
    line_y = dot_y + dot_r - Inches(0.02)

    # Horizontal connector line
    _rect(sl, Inches(0.5) + step_w * 0.5, line_y,
          Inches(12.5) - step_w, Pt(2), _rgb(0x37, 0x41, 0x51))

    for i, step in enumerate(steps):
        ac = _accent_rgb(i)
        cx = Inches(0.5) + i * step_w + step_w / 2

        # Circle
        circle = sl.shapes.add_shape(
            9,  # MSO_SHAPE_TYPE.OVAL
            cx - dot_r, dot_y, dot_r * 2, dot_r * 2,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ac
        circle.line.fill.background()

        # Number inside circle
        tf = circle.text_frame
        tf.word_wrap = False
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run = p2.add_run()
        run.text = str(step.get("num", str(i + 1).zfill(2)))
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = _rgb(0xFF, 0xFF, 0xFF)

        lx = cx - step_w * 0.45
        lw = step_w * 0.9

        _txt(sl, _e_plain(step.get("label", "")), lx, dot_y + dot_r * 2 + Inches(0.18),
             lw, Inches(0.5), size=11, bold=True, color=_DARK_TEXT, align=PP_ALIGN.CENTER)
        if step.get("sub"):
            _txt(sl, _e_plain(step["sub"]), lx, dot_y + dot_r * 2 + Inches(0.72),
                 lw, Inches(0.45), size=9, color=_MUTED, align=PP_ALIGN.CENTER)
        if step.get("duration"):
            _txt(sl, _e_plain(step["duration"]), lx, dot_y + dot_r * 2 + Inches(1.2),
                 lw, Inches(0.35), size=8, bold=True, color=_rgb(0x94, 0xa3, 0xb8),
                 align=PP_ALIGN.CENTER)


def _pptx_split3(prs, s: dict, idx: int):
    sl = _blank(prs)
    _set_bg(sl, _DARK_BG)
    ac = _accent_rgb(idx)
    _accent_bar(sl, ac)

    _txt(sl, s.get("title", ""), Inches(0.7), Inches(0.3), Inches(11), Inches(0.8),
         size=24, bold=True, color=_DARK_TEXT)

    cols_data = s.get("columns", [])[:3]
    col_w = Inches(4.0)
    for i, col in enumerate(cols_data):
        col_ac = _rgb(*_ACCENTS_RGB[i % len(_ACCENTS_RGB)])
        x = Inches(0.25) + i * (col_w + Inches(0.2))
        _rect(sl, x, Inches(1.3), col_w, Inches(5.9), _rgb(0x14, 0x1d, 0x2e),
              line_color=_rgb(0x28, 0x30, 0x48))
        _rect(sl, x, Inches(1.3), col_w, Inches(0.55), col_ac)
        _txt(sl, _e_plain(col.get("heading", "")), x, Inches(1.3), col_w, Inches(0.55),
             size=11, bold=True,
             color=_LIGHT_TEXT if i == 1 else _DARK_TEXT,
             align=PP_ALIGN.CENTER)
        cy = Inches(2.0)
        if col.get("body"):
            _txt(sl, _e_plain(col["body"]), x + Inches(0.15), cy, col_w - Inches(0.3), Inches(1.2),
                 size=10, color=_MUTED)
            cy += Inches(1.25)
        for b in col.get("bullets", [])[:5]:
            _txt(sl, f"• {_e_plain(b)}", x + Inches(0.15), cy, col_w - Inches(0.3), Inches(0.5),
                 size=10, color=_DARK_TEXT)
            cy += Inches(0.52)


_RENDERERS = {
    "title":    _pptx_title,
    "agenda":   _pptx_agenda,
    "timeline": _pptx_timeline,
    "section":  _pptx_section,
    "content":  _pptx_content,
    "split":    _pptx_split,
    "split3":   _pptx_split3,
    "stats":    _pptx_stats,
    "table":    _pptx_table,
    "cases":    _pptx_cases,
    "quote":    _pptx_quote,
    "closing":  _pptx_closing,
}


def render_pptx(title: str, slides: list[dict], out_path: Path) -> bool:
    if not _PPTX_OK:
        return False
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    for i, slide in enumerate(slides):
        renderer = _RENDERERS.get(slide.get("type", "content"), _pptx_content)
        try:
            renderer(prs, slide, i)
        except Exception:
            pass
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return True
